#!/usr/bin/env python3
"""
Star Factor Pitch Deck Generator
Creates a PPTX from PITCH-DECK-V2.md with properly centered content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import re

def parse_markdown_file(filepath):
    """Parse markdown file into slide sections"""
    with open(filepath, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by slide markers
    slides = []
    current_slide = {'title': '', 'content': []}

    lines = content.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i].strip()

        # Skip empty lines at start
        if not line and not current_slide['content']:
            i += 1
            continue

        # Slide separator
        if line == '---':
            if current_slide['title'] or current_slide['content']:
                slides.append(current_slide)
                current_slide = {'title': '', 'content': []}
            i += 1
            continue

        # Slide title (## SLIDE X:)
        if line.startswith('## SLIDE'):
            match = re.match(r'## SLIDE \d+: (.+)', line)
            if match:
                current_slide['title'] = match.group(1)
            i += 1
            continue

        # Main title slides (single # or ##)
        if line.startswith('# ') and not line.startswith('## '):
            current_slide['title'] = line.lstrip('# ').strip()
            i += 1
            continue

        if line.startswith('## ') and not line.startswith('## SLIDE'):
            current_slide['title'] = line.lstrip('# ').strip()
            i += 1
            continue

        # Content
        if line:
            current_slide['content'].append(line)

        i += 1

    # Add last slide
    if current_slide['title'] or current_slide['content']:
        slides.append(current_slide)

    return slides

def add_centered_title_slide(prs, title, subtitle=''):
    """Add a title slide with centered content"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(9), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    return slide

def add_content_slide(prs, title, content_lines):
    """Add a content slide with properly centered elements"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title at top
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Content area - centered vertically
    content_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.5), Inches(8.5), Inches(5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Center vertically

    # Parse and add content
    is_table = False
    table_lines = []
    code_block = False
    code_lines = []

    for line in content_lines:
        # Handle code blocks (for diagrams)
        if line.startswith('```'):
            if code_block:
                # End code block
                add_code_block(content_frame, code_lines)
                code_lines = []
                code_block = False
            else:
                code_block = True
            continue

        if code_block:
            code_lines.append(line)
            continue

        # Handle tables
        if line.startswith('|'):
            is_table = True
            table_lines.append(line)
            continue
        elif is_table and not line.startswith('|'):
            add_table_to_frame(content_frame, table_lines)
            table_lines = []
            is_table = False

        # Regular content
        if line.startswith('**') and line.endswith('**'):
            # Bold header
            p = content_frame.add_paragraph()
            p.text = line.strip('*')
            p.font.size = Pt(18)
            p.font.bold = True
            p.space_before = Pt(12)
            p.alignment = PP_ALIGN.CENTER
        elif line.startswith('*') and not line.startswith('**'):
            # Italic
            p = content_frame.add_paragraph()
            p.text = line.strip('*')
            p.font.size = Pt(14)
            p.font.italic = True
            p.alignment = PP_ALIGN.CENTER
        elif line.startswith('- ') or line.startswith('• '):
            # Bullet point
            p = content_frame.add_paragraph()
            p.text = line.lstrip('- •').strip()
            p.level = 0
            p.font.size = Pt(14)
            p.space_before = Pt(6)
        else:
            # Regular text
            p = content_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(14)
            p.space_before = Pt(6)

    # Handle remaining table or code block
    if table_lines:
        add_table_to_frame(content_frame, table_lines)
    if code_lines:
        add_code_block(content_frame, code_lines)

    return slide

def add_table_to_frame(frame, table_lines):
    """Add table content to text frame"""
    for line in table_lines:
        if line.startswith('|---') or line.startswith('| ---'):
            continue  # Skip separator

        cells = [c.strip() for c in line.split('|')[1:-1]]
        p = frame.add_paragraph()
        p.text = '  •  '.join(cells)
        p.font.size = Pt(12)
        p.space_before = Pt(4)

def add_code_block(frame, code_lines):
    """Add code block (diagram) to frame"""
    for line in code_lines:
        p = frame.add_paragraph()
        p.text = line
        p.font.name = 'Courier New'
        p.font.size = Pt(10)
        p.space_before = Pt(2)

def create_pitch_deck():
    """Main function to create the pitch deck"""
    # Load markdown
    slides_data = parse_markdown_file('../PITCH-DECK-V2.md')

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Process slides
    for i, slide_data in enumerate(slides_data):
        title = slide_data['title']
        content = slide_data['content']

        # First slide special handling
        if i == 0 and title == 'STAR FACTOR':
            subtitle = content[0] if content else ''
            add_centered_title_slide(prs, title, subtitle)
        # Contact/closing slides
        elif 'bolaji@chainfren.com' in '\n'.join(content).lower() or title.startswith('Ready to'):
            # Simple centered slide
            add_centered_title_slide(prs, title, '\n'.join(content[:3]))
        else:
            # Regular content slide
            add_content_slide(prs, title, content)

    # Save
    output_file = 'Star-Factor-Pitch-Deck-v2.pptx'
    prs.save(output_file)
    print(f'✅ Created: {output_file}')
    print(f'   Slides: {len(prs.slides)}')
    print(f'   All content centered and fully editable')

if __name__ == '__main__':
    create_pitch_deck()
