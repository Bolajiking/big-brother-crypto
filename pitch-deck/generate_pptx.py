#!/usr/bin/env python3
"""
Star Factor Pitch Deck Generator
Creates PPTX from HTML design - 12 slides, properly centered
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Design system colors from HTML
BG = RGBColor(9, 1, 27)  # #09011B
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(140, 140, 140)  # rgba(255,255,255,0.55)
BLUE = RGBColor(64, 172, 255)  # #40ACFF
CYAN = RGBColor(90, 205, 255)  # #5ACDFF
MINT = RGBColor(203, 240, 184)  # #CBF0B8
PERIWINKLE = RGBColor(141, 170, 255)  # #8DAAFF

def create_presentation():
    """Create base presentation with 16:9 dimensions"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 ratio
    return prs

def add_footer(slide, prs):
    """Add footer to slide"""
    footer_y = Inches(5.1)
    # Left: Logo placeholder
    logo_box = slide.shapes.add_textbox(Inches(0.5), footer_y, Inches(1), Inches(0.3))
    logo_frame = logo_box.text_frame
    logo_frame.text = "⭐"
    logo_frame.paragraphs[0].font.size = Pt(14)
    logo_frame.paragraphs[0].font.color.rgb = MUTED

    # Right: Tag
    tag_box = slide.shapes.add_textbox(Inches(7), footer_y, Inches(2.5), Inches(0.3))
    tag_frame = tag_box.text_frame
    tag_frame.text = "STAR FACTOR — PRE-SEED 2026"
    tag_frame.paragraphs[0].font.size = Pt(9)
    tag_frame.paragraphs[0].font.color.rgb = MUTED
    tag_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

def add_section_label(slide, text, top=0.5):
    """Add section label (blue uppercase)"""
    label_box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(0.3))
    label_frame = label_box.text_frame
    label_frame.text = text.upper()
    p = label_frame.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.LEFT

def add_centered_text(slide, text, top, height, font_size, bold=False, color=WHITE, align=PP_ALIGN.CENTER):
    """Add centered text box"""
    box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(9), Inches(height))
    frame = box.text_frame
    frame.text = text
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box

# SLIDE 1: Title
def create_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Logo row
    add_centered_text(slide, "A CHAINFREN PRODUCTION", 1.2, 0.3, 10, color=MUTED)

    # Main title
    title_box = add_centered_text(slide, "STAR FACTOR", 1.8, 0.8, 54, bold=True, color=CYAN)

    # Tagline
    add_centered_text(slide, "WATCH · PREDICT · EARN", 2.7, 0.4, 16, color=MUTED)

    # Meta row
    meta_text = "Pre-Seed Round  •  $200,000  •  April 2026"
    add_centered_text(slide, meta_text, 3.5, 0.4, 14, color=MUTED)

# SLIDE 2: The Problem
def create_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_label(slide, "The Problem")

    # Title
    title_box = add_centered_text(slide, "Fans create all the value.\nFans keep none of it.", 0.9, 0.7, 32, bold=True)

    # Stat callout
    add_centered_text(slide, "2B+", 0.9, 0.5, 42, bold=True, color=CYAN, align=PP_ALIGN.RIGHT)
    stat_sub = slide.shapes.add_textbox(Inches(7), Inches(1.4), Inches(2.5), Inches(0.5))
    stat_sub.text_frame.text = "votes cast per BBNaija season"
    stat_sub.text_frame.paragraphs[0].font.size = Pt(10)
    stat_sub.text_frame.paragraphs[0].font.color.rgb = MUTED
    stat_sub.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Quote
    quote_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(0.6))
    quote_frame = quote_box.text_frame
    quote_frame.text = '"I watched every episode, voted 200 times, spent ₦50,000... and I got nothing."'
    quote_frame.paragraphs[0].font.size = Pt(13)
    quote_frame.paragraphs[0].font.italic = True
    quote_frame.paragraphs[0].font.color.rgb = MUTED

    # Two columns
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(4.3), Inches(2))
    left_frame = left_box.text_frame
    left_frame.text = "What Fans Give\n• 100+ hours per season\n• ₦10K–50K+ on SMS votes\n• Social campaigns that build stars\n• Cultural influence driving brand deals"
    for p in left_frame.paragraphs:
        p.font.size = Pt(11)
        p.space_before = Pt(4)
    left_frame.paragraphs[0].font.bold = True
    left_frame.paragraphs[0].font.color.rgb = RGBColor(255, 107, 107)

    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.7), Inches(4.3), Inches(2))
    right_frame = right_box.text_frame
    right_frame.text = "What Fans Get\n• Zero financial return\n• No ownership, no stake\n• No share of post-show revenue\n• A 'thank you'"
    for p in right_frame.paragraphs:
        p.font.size = Pt(11)
        p.space_before = Pt(4)
    right_frame.paragraphs[0].font.bold = True
    right_frame.paragraphs[0].font.color.rgb = BLUE

    add_footer(slide, prs)

# SLIDE 3: The Market
def create_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_label(slide, "The Market")

    add_centered_text(slide, "Two $3B+ markets. One intersection.\nZero competitors.", 0.9, 0.6, 28, bold=True)
    add_centered_text(slide, "$7.6B total addressable market at the convergence of entertainment and prediction.", 1.6, 0.3, 14, color=MUTED)

    # Four stats
    stats = [
        ("$4B+", "Nigerian Entertainment\nIndustry (15% YoY)", CYAN),
        ("$3.6B", "Nigerian Gross Gaming\nRevenue", BLUE),
        ("50M+", "Reality TV viewers\nper BBNaija season", PERIWINKLE),
        ("60M+", "Active Nigerian\nbettors", MINT)
    ]

    for i, (num, label, color) in enumerate(stats):
        x = 0.5 + (i * 2.3)
        stat_box = slide.shapes.add_textbox(Inches(x), Inches(2.1), Inches(2.2), Inches(1.2))
        frame = stat_box.text_frame
        frame.text = f"{num}\n{label}"
        frame.paragraphs[0].font.size = Pt(24)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[0].font.color.rgb = color
        for p in frame.paragraphs[1:]:
            p.font.size = Pt(10)
            p.font.color.rgb = MUTED

    # Insight box
    insight = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    insight.text_frame.text = "📈 In-play betting grew from 25% to 40% of all bets in two years. Nigerians want to predict outcomes while watching live events. We give them a better event to predict on."
    insight.text_frame.paragraphs[0].font.size = Pt(12)
    insight.text_frame.paragraphs[0].font.color.rgb = WHITE

    add_footer(slide, prs)

# SLIDE 4: Why Now
def create_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_label(slide, "Timing")

    add_centered_text(slide, "Five forces converging in 2026", 0.9, 0.5, 28, bold=True)

    forces = [
        ("01", "Mobile video crossed the threshold", "140M smartphones. Data usage up 93% in 2 years. 24/7 streaming is now viable."),
        ("02", "Prediction markets are the 2026 breakout category", "$13B+ monthly trading. a16z, Galaxy Research, Pantera all published major theses in the last 6 months."),
        ("03", "Decentralized streaming slashed production costs", "Livepeer cuts costs 50–80% vs AWS. Fishtank.live proved it: $3M revenue, $486K profit."),
        ("04", "Nigerian crypto rails are mature", "Paystack ($10B+ processing). cNGN launched 2025. Infrastructure is production-ready."),
        ("05", "BBNaija off-season creates a content vacuum", "Season 9 ends October. November–May: no major reality content. We launch into maximum attention gap.")
    ]

    y_start = 1.5
    for i, (num, title, desc) in enumerate(forces):
        if i < 4:
            row = i // 2
            col = i % 2
            x = 0.5 + (col * 4.7)
            y = y_start + (row * 0.85)
            h = 0.75
        else:
            x = 0.5
            y = y_start + 1.7
            h = 0.6

        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(4.5 if i < 4 else 9), Inches(h))
        frame = box.text_frame
        frame.text = f"{num}  {title}\n{desc}"
        frame.paragraphs[0].font.size = Pt(12)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[1].font.size = Pt(10)
        frame.paragraphs[1].font.color.rgb = MUTED

    add_footer(slide, prs)

# SLIDE 5: The Product
def create_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_label(slide, "The Product")

    add_centered_text(slide, "24/7 reality TV where fans have\nreal skin in the game.", 0.9, 0.6, 28, bold=True)
    add_centered_text(slide, "Watch · Predict · Earn — three layers, one flywheel.", 1.5, 0.3, 13, color=MUTED)

    pillars = [
        ("01 WATCH", "8 cameras. Every room. 24/7. Mobile-first. Powered by Livepeer decentralized streaming.", "50–80% cheaper than AWS · HLS playback", PERIWINKLE),
        ("02 PREDICT", "Who gets evicted? Who wins Head of House? Fans place real-money predictions. Right calls earn real money.", "5% platform rake · On-chain settlement via Solana", BLUE),
        ("03 EARN", "Contestants earn from fan tips during the show — not just after. Top supporters get exclusive access.", "30% platform / 70% contestant split · Instant cashout", MINT)
    ]

    for i, (title, desc, detail, color) in enumerate(pillars):
        x = 0.5 + (i * 3.2)
        box = slide.shapes.add_textbox(Inches(x), Inches(1.9), Inches(3), Inches(2.2))
        frame = box.text_frame
        frame.text = f"{title}\n\n{desc}\n\n{detail}"
        frame.paragraphs[0].font.size = Pt(16)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[0].font.color.rgb = color
        frame.paragraphs[1].font.size = Pt(11)
        frame.paragraphs[2].font.size = Pt(9)
        frame.paragraphs[2].font.color.rgb = MUTED

    add_centered_text(slide, "Predictions → more watching → better predictions → winning → dopamine → repeat", 4.4, 0.3, 10, color=MUTED)

    add_footer(slide, prs)

# SLIDE 6: Business Model
def create_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_label(slide, "Business Model")

    add_centered_text(slide, "Nine revenue streams.", 0.9, 0.4, 28, bold=True)
    add_centered_text(slide, "Break-even on sponsorship alone. User revenue is pure upside.", 1.3, 0.3, 13, color=MUTED)

    # Revenue streams (left side)
    streams = [
        ("Sponsorship & Brands", "25%", "Title sponsor + room, challenge, digital packages"),
        ("Stakes Sales", "35%", "Naira → Stakes · 15–20% margin"),
        ("Subscriptions", "15%", "4 tiers: Bronze ₦480/wk → Platinum ₦160K/season"),
        ("Prediction Market Fees", "15%", "5% rake on all prediction volume"),
        ("Interactive (TTS/Gifts)", "5%", "₦480–1,600 per interaction"),
        ("Contestant Tips + Content", "5%", "30% platform take · YouTube highlights")
    ]

    y = 1.7
    for name, pct, desc in streams:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(5), Inches(0.4))
        frame = box.text_frame
        frame.text = f"{name} ({pct}) — {desc}"
        frame.paragraphs[0].font.size = Pt(9)
        y += 0.45

    # Currencies (right side)
    curr_box1 = slide.shapes.add_textbox(Inches(5.8), Inches(1.7), Inches(3.7), Inches(1.3))
    frame1 = curr_box1.text_frame
    frame1.text = "CLOUT (Free Currency)\nEarned by watching, chatting, logging in. Creates habit. Creates the audience sponsors pay for."
    frame1.paragraphs[0].font.size = Pt(14)
    frame1.paragraphs[0].font.bold = True
    frame1.paragraphs[0].font.color.rgb = PERIWINKLE
    frame1.paragraphs[1].font.size = Pt(10)
    frame1.paragraphs[1].font.color.rgb = MUTED

    curr_box2 = slide.shapes.add_textbox(Inches(5.8), Inches(3.2), Inches(3.7), Inches(1.3))
    frame2 = curr_box2.text_frame
    frame2.text = "STAKES (Premium Currency)\nPurchased with Naira (₦200 min). Pegged at $0.01 USD. Used for predictions, premium votes, TTS. Cashable to bank."
    frame2.paragraphs[0].font.size = Pt(14)
    frame2.paragraphs[0].font.bold = True
    frame2.paragraphs[0].font.color.rgb = CYAN
    frame2.paragraphs[1].font.size = Pt(10)
    frame2.paragraphs[1].font.color.rgb = MUTED

    add_footer(slide, prs)

# SLIDE 7: Projections
def create_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_label(slide, "Unit Economics & Season 1 Projections")

    add_centered_text(slide, "Profitable at base case. Break-even at conservative.", 0.9, 0.4, 26, bold=True)

    # Table
    table_data = [
        ("", "Conservative", "Base Case", "Stretch"),
        ("Registered Users", "15,000", "40,000", "100,000"),
        ("Paying Users", "1,200 (8%)", "4,000 (10%)", "12,000 (12%)"),
        ("User Revenue", "₦15M (~$9K)", "₦60M (~$37K)", "₦200M (~$125K)"),
        ("Sponsorship", "₦15M (~$9K)", "₦40M (~$25K)", "₦80M (~$50K)"),
        ("Production Cost", "₦20M (~$13K) fixed", "", ""),
        ("Total Gross Revenue", "Break-even", "Profitable", "Highly Profitable")
    ]

    y_start = 1.4
    for i, row in enumerate(table_data):
        for j, cell in enumerate(row):
            if cell:
                x = 0.5 + (j * 2.3)
                box = slide.shapes.add_textbox(Inches(x), Inches(y_start + i * 0.35), Inches(2.2), Inches(0.3))
                frame = box.text_frame
                frame.text = cell
                p = frame.paragraphs[0]
                p.font.size = Pt(9 if i > 0 else 10)
                p.font.bold = (i == 0 or i == len(table_data) - 1)
                p.font.color.rgb = MUTED if j == 0 else WHITE

    # Metrics (right side)
    metrics = [
        ("CAC (Organic)", "₦100–300"),
        ("Free → Paid Conversion", "8–10%"),
        ("ARPPU per Season", "₦8K–12K"),
        ("LTV / CAC (Organic)", "5–8×")
    ]

    y = 1.4
    for label, val in metrics:
        box = slide.shapes.add_textbox(Inches(6.5), Inches(y), Inches(3), Inches(0.6))
        frame = box.text_frame
        frame.text = f"{label}\n{val}"
        frame.paragraphs[0].font.size = Pt(9)
        frame.paragraphs[0].font.color.rgb = MUTED
        frame.paragraphs[1].font.size = Pt(14)
        frame.paragraphs[1].font.bold = True
        y += 0.75

    add_footer(slide, prs)

# SLIDE 8: Competitive Landscape
def create_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_label(slide, "Competitive Landscape")

    add_centered_text(slide, "Nobody occupies our intersection.", 0.9, 0.5, 26, bold=True)

    # Comparison matrix
    competitors = ["Star Factor", "BBNaija", "BetKing", "Fishtank", "Polymarket"]
    features = [
        "Owned Live Content",
        "Prediction Markets",
        "Fan Revenue Share",
        "African Market Native",
        "On-chain Settlement"
    ]

    # Star Factor has all checkmarks
    matrix = [
        [True, True, False, True, False],  # Owned Live Content
        [True, False, True, False, True],  # Prediction Markets
        [True, False, True, False, False], # Fan Revenue Share
        [True, True, True, False, False],  # African Market Native
        [True, False, False, False, True]  # On-chain Settlement
    ]

    # Headers
    y_header = 1.5
    for i, comp in enumerate(competitors):
        x = 2.5 + (i * 1.4)
        box = slide.shapes.add_textbox(Inches(x), Inches(y_header), Inches(1.3), Inches(0.3))
        frame = box.text_frame
        frame.text = comp
        p = frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = CYAN if i == 0 else MUTED
        p.alignment = PP_ALIGN.CENTER

    # Feature rows
    for i, feature in enumerate(features):
        y = 1.9 + (i * 0.45)
        # Feature name
        feat_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(2), Inches(0.4))
        feat_box.text_frame.text = feature
        feat_box.text_frame.paragraphs[0].font.size = Pt(10)

        # Checkmarks/crosses for each competitor
        for j in range(5):
            x = 2.5 + (j * 1.4) + 0.5
            check_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.3), Inches(0.3))
            check_box.text_frame.text = "✓" if matrix[i][j] else "✗"
            p = check_box.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = CYAN if (j == 0 and matrix[i][j]) else MUTED
            p.alignment = PP_ALIGN.CENTER

    add_centered_text(slide, "★ Only Star Factor has all five.", 4.2, 0.3, 12, bold=True, color=CYAN)

    add_footer(slide, prs)

# SLIDE 9: Traction
def create_slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_label(slide, "Traction")

    add_centered_text(slide, "Built, not planned.", 0.9, 0.4, 28, bold=True)
    add_centered_text(slide, "We are not asking for money to build a prototype — we are asking for money to turn on the cameras.", 1.3, 0.4, 12, color=MUTED)

    # Traction items
    items = [
        ("LIVE", "Streaming Platform", "8-camera multi-feed via Livepeer · HLS playback · Mobile-first UI", MINT),
        ("LIVE", "Real-Time Chat", "Socket.io · Emoji reactions · SFX triggers · TTS interface", MINT),
        ("LIVE", "Web3 Auth & Wallet", "Privy · Email/Google → embedded Solana wallet", MINT),
        ("LIVE", "Payment Rails", "Paystack · Naira deposits · Instant bank withdrawals", MINT),
        ("LIVE", "Contestant Applications", "6-step application form live · Collecting Season 1 applicants", MINT),
        ("IN PROGRESS", "Prediction Markets", "Betting UI, odds engine, settlement · Shipping in 4 weeks", BLUE)
    ]

    y = 1.8
    for status, name, desc, color in items:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(4.5), Inches(0.55))
        frame = box.text_frame
        frame.text = f"[{status}] {name}\n{desc}"
        frame.paragraphs[0].font.size = Pt(11)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[1].font.size = Pt(9)
        frame.paragraphs[1].font.color.rgb = MUTED
        y += 0.6

    # Validation stats
    stats_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4.2), Inches(2))
    stats_frame = stats_box.text_frame
    stats_frame.text = "Validation:\n\n3M+ social views — Fanplay (Solana)\n\n3/4 Colosseum Consumer Apps 1st-place winners are prediction plays\n\n$486K — Fishtank.live profit on $3M revenue using same Livepeer stack"
    for i, p in enumerate(stats_frame.paragraphs):
        p.font.size = Pt(11 if i == 0 else 10)
        p.font.bold = (i == 0)
        p.space_before = Pt(6)

    add_footer(slide, prs)

# SLIDE 10: Team
def create_slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_label(slide, "The Team")

    add_centered_text(slide, "Built by Chainfren.", 0.9, 0.4, 28, bold=True)

    team = [
        ("Bolaji Majiyagbe", "Founder & Product Lead", "5 years building in Web3. Decentralized infrastructure specialist. Deep understanding of Nigerian consumer behavior, payment rails, and digital entertainment."),
        ("Bashir Jaji", "Lead Engineer", "Full-stack developer. Built the entire Star Factor platform from scratch — streaming, chat, wallet, payments, Web3 auth."),
        ("Kunle Aikulola", "Marketing & Growth Lead", "MSc in Marketing. Nigerian digital audience specialist. Influencer network spanning entertainment, pop culture, and sports verticals.")
    ]

    for i, (name, role, bio) in enumerate(team):
        x = 0.5 + (i * 3.2)
        box = slide.shapes.add_textbox(Inches(x), Inches(1.7), Inches(3), Inches(2))
        frame = box.text_frame
        frame.text = f"{name}\n{role}\n\n{bio}"
        frame.paragraphs[0].font.size = Pt(13)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[1].font.size = Pt(10)
        frame.paragraphs[1].font.color.rgb = BLUE
        frame.paragraphs[2].font.size = Pt(9)
        frame.paragraphs[2].font.color.rgb = MUTED

    # Hiring
    hiring_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.5))
    hiring_frame = hiring_box.text_frame
    hiring_frame.text = "Hiring with this raise: Show Runner · Content Editor · Community Manager + Legal counsel · Nigerian entertainment advisor"
    hiring_frame.paragraphs[0].font.size = Pt(11)
    hiring_frame.paragraphs[0].font.color.rgb = MUTED

    add_footer(slide, prs)

# SLIDE 11: The Ask
def create_slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_label(slide, "The Ask")

    # Main ask
    ask_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(4), Inches(1.5))
    ask_frame = ask_box.text_frame
    ask_frame.text = "Pre-Seed Round\n$200K\nTo complete the platform and launch Season 1 in Lagos, Q4 2026."
    ask_frame.paragraphs[0].font.size = Pt(14)
    ask_frame.paragraphs[0].font.color.rgb = MUTED
    ask_frame.paragraphs[1].font.size = Pt(44)
    ask_frame.paragraphs[1].font.bold = True
    ask_frame.paragraphs[1].font.color.rgb = CYAN
    ask_frame.paragraphs[2].font.size = Pt(12)
    ask_frame.paragraphs[2].font.color.rgb = MUTED

    # Parameters
    params = [
        ("Season Duration", "6 weeks · 42 days live"),
        ("Contestants", "16–20 participants"),
        ("Location", "Lagos, Nigeria"),
        ("Launch Window", "Q4 2026 (BBNaija off-season)"),
        ("Format", "24/7 live · weekly evictions · daily challenges"),
        ("Series A Trigger", "$500K–1M seed post Season 1")
    ]

    y = 0.9
    for label, val in params:
        box = slide.shapes.add_textbox(Inches(5), Inches(y), Inches(4.5), Inches(0.35))
        frame = box.text_frame
        frame.text = f"{label}: {val}"
        frame.paragraphs[0].font.size = Pt(10)
        y += 0.38

    # Fund allocation
    funds = [
        ("35%", "$70K", "Operations & Production", "House rental · cameras · contestant stipends · production staff"),
        ("25%", "$50K", "Product & Engineering", "Complete prediction markets · admin dashboard · analytics · QA · hosting"),
        ("25%", "$50K", "Marketing & Launch", "30 micro-influencer partnerships · TikTok/YouTube content · WhatsApp campaigns"),
        ("15%", "$30K", "Prize Pool", "Grand prize ($15K) · weekly challenge prizes · contestant appearance fees")
    ]

    for i, (pct, amt, cat, detail) in enumerate(funds):
        x = 0.5 + (i * 2.4)
        box = slide.shapes.add_textbox(Inches(x), Inches(3), Inches(2.3), Inches(1.5))
        frame = box.text_frame
        frame.text = f"{pct}\n{amt}\n{cat}\n{detail}"
        frame.paragraphs[0].font.size = Pt(18)
        frame.paragraphs[0].font.bold = True
        frame.paragraphs[0].font.color.rgb = CYAN
        frame.paragraphs[1].font.size = Pt(14)
        frame.paragraphs[1].font.bold = True
        frame.paragraphs[2].font.size = Pt(11)
        frame.paragraphs[2].font.bold = True
        frame.paragraphs[3].font.size = Pt(9)
        frame.paragraphs[3].font.color.rgb = MUTED

    add_footer(slide, prs)

# SLIDE 12: Close
def create_slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_centered_text(slide, "Star Factor · Pre-Seed 2026", 1, 0.3, 12, color=MUTED)
    add_centered_text(slide, "Ready to discuss?", 1.5, 0.6, 38, bold=True, color=CYAN)

    # Milestones
    milestones = [
        "✓ Platform is built",
        "✓ Payments work",
        "✓ Streaming works",
        "✓ Applications open",
        "→ Prediction markets in 4 weeks",
        "→ Cameras on Q4 2026"
    ]

    y = 2.3
    for m in milestones:
        add_centered_text(slide, m, y, 0.25, 11, color=MINT if m.startswith("✓") else BLUE)
        y += 0.28

    # Contact
    add_centered_text(slide, "Bolaji Majiyagbe", 4.2, 0.25, 16, bold=True)
    add_centered_text(slide, "bolaji@chainfren.com", 4.5, 0.25, 14, color=CYAN)

    # Footer
    add_centered_text(slide, "⭐ Built by Chainfren · A Web3 creative studio building for the creator economy", 5, 0.3, 9, color=MUTED)

def main():
    """Generate the complete pitch deck"""
    print("Generating Star Factor Pitch Deck...")

    prs = create_presentation()

    # Create all 12 slides
    create_slide_1(prs)
    create_slide_2(prs)
    create_slide_3(prs)
    create_slide_4(prs)
    create_slide_5(prs)
    create_slide_6(prs)
    create_slide_7(prs)
    create_slide_8(prs)
    create_slide_9(prs)
    create_slide_10(prs)
    create_slide_11(prs)
    create_slide_12(prs)

    # Save
    output_file = 'Star-Factor-Pitch-Deck.pptx'
    prs.save(output_file)

    print(f"✅ Created: {output_file}")
    print(f"   12 slides, properly centered, fully editable")
    print(f"   Based on HTML design from claude.ai/design")

if __name__ == '__main__':
    main()
